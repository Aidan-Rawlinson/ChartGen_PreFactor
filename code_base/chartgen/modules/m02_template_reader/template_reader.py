"""
template_reader.py
M02 — Template Reader

Reads a .pptx template file and produces:
  1. A list of chart placeholders (name, position, size) per slide
  2. Any yellow URL textboxes that are fully contained within those placeholders
  3. A cleaned copy of the template with yellow textboxes stripped

Yellow textbox detection criteria (all three must be true):
  - Fill colour in HSV: hue 40°–70°, saturation > 50%, value > 50%
  - The textbox bounds are fully inside a chart placeholder on the same slide
  - The textbox text contains a URL matching the toolkit URL pattern

One URL per placeholder is the contract. If multiple qualifying textboxes
fall within the same placeholder, a warning is raised and the first is used.
"""

import re
import copy
import os
from colorsys import rgb_to_hsv
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# Hue range for "human yellow" in degrees
YELLOW_HUE_MIN = 40
YELLOW_HUE_MAX = 70
YELLOW_SAT_MIN = 0.50   # saturation floor — rejects cream / off-white
YELLOW_VAL_MIN = 0.50   # value floor — rejects dark mustard

# Toolkit URL pattern — matches NHS Benchmarking members URLs
TOOLKIT_URL_RE = re.compile(
    r"https?://[^\s\"'<>]+"
    r"nhsbenchmarking[^\s\"'<>]*",
    re.IGNORECASE,
)

# Fallback: any http/https URL
GENERIC_URL_RE = re.compile(r"https?://[^\s\"'<>]+", re.IGNORECASE)

# Image file extensions recognised by insert_picture
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".wmf", ".emf"}

# Excel workbook pattern: full path ending in .xlsx or .xlsm,
# followed by [key:rangename,...] bracket block
# e.g. C:\Projects\analysis.xlsx[driver:UNIT,export:SummaryTable]
EXCEL_PATH_RE = re.compile(
    r"^(.+\.(?:xlsx|xlsm))\[([^\]]+)\]$",
    re.IGNORECASE,
)

# ---------------------------------------------------------------------------
# Data structures
# ---------------------------------------------------------------------------

@dataclass
class PlaceholderInfo:
    """Describes a single named placeholder on a slide."""
    slide_index: int          # 0-based
    name: str                 # e.g. "Chart 1"
    left: int                 # EMU
    top: int                  # EMU
    width: int                # EMU
    height: int               # EMU
    # Yellow box classification — one of: "chart", "picture", "excel", ""
    content_type: str = ""
    # chart: toolkit URL
    url: str = ""
    label: str = ""           # user text from the yellow textbox (for urls.csv)
    # picture: image file path (may contain [code] or [id] tokens)
    image_path: str = ""
    # excel: workbook path, export range name, driver range name (optional)
    excel_path: str = ""
    excel_export_range: str = ""
    excel_driver_range: str = ""


@dataclass
class TemplateReadResult:
    """Full output of a template read operation."""
    slide_width: int          # EMU
    slide_height: int         # EMU
    placeholders: list = field(default_factory=list)   # list[PlaceholderInfo]
    warnings: list = field(default_factory=list)       # list[str]
    cleaned_pptx_bytes: bytes = b""                    # template with yellow boxes removed


# ---------------------------------------------------------------------------
# Colour helpers
# ---------------------------------------------------------------------------

def _rgb_to_hsv_degrees(r, g, b):
    """Convert 0–255 RGB to (hue_degrees, sat_0_1, val_0_1)."""
    h, s, v = rgb_to_hsv(r / 255, g / 255, b / 255)
    return h * 360, s, v


def _is_yellow(rgb) -> bool:
    """
    Return True if the RGB colour is a human-visible yellow.
    RGBColor is a tuple subclass — access via index, not .red/.green/.blue.
    """
    if rgb is None:
        return False
    try:
        r, g, b = rgb[0], rgb[1], rgb[2]
    except (TypeError, IndexError):
        return False
    h, s, v = _rgb_to_hsv_degrees(r, g, b)
    return (
        YELLOW_HUE_MIN <= h <= YELLOW_HUE_MAX
        and s >= YELLOW_SAT_MIN
        and v >= YELLOW_VAL_MIN
    )


def _get_shape_fill_rgb(shape):
    """
    Return the solid fill RGBColor of a shape, or None if not determinable.

    python-pptx's fore_color.rgb raises when the colour is stored as a theme
    colour reference (even if an srgbClr override exists in the XML). We read
    the raw XML directly to catch both cases reliably.
    """
    import re as _re
    try:
        xml = shape._element.xml
        # Look for explicit sRGB values anywhere in the shape XML
        matches = _re.findall(r'<a:srgbClr\s+val="([0-9A-Fa-f]{6})"', xml)
        if matches:
            hex_val = matches[0].upper()
            r = int(hex_val[0:2], 16)
            g = int(hex_val[2:4], 16)
            b = int(hex_val[4:6], 16)
            from pptx.dml.color import RGBColor
            return RGBColor(r, g, b)
    except Exception:
        pass
    # Fallback: try python-pptx native accessor
    try:
        fill = shape.fill
        if str(fill.type) == "SOLID (1)":
            return fill.fore_color.rgb
    except Exception:
        pass
    return None


def _extract_url_from_text(text: str) -> str:
    """
    Extract the first toolkit URL from text.
    Falls back to any HTTP URL if no toolkit URL is found.
    Returns empty string if nothing found.
    """
    m = TOOLKIT_URL_RE.search(text)
    if m:
        return m.group(0).rstrip(".,;)")
    m = GENERIC_URL_RE.search(text)
    if m:
        return m.group(0).rstrip(".,;)")
    return ""


def _classify_yellow_box(text: str) -> dict:
    """
    Classify the content of a yellow textbox and return a dict describing it.

    Returns a dict with key 'content_type' set to one of:
      "chart"   — text contains a toolkit or HTTP URL
      "picture" — text is a path to a supported image file
      "excel"   — text matches the Excel path + range bracket syntax
      ""        — unrecognised; yellow box will be stripped but ignored

    Additional keys depend on content_type:
      chart:   url, label
      picture: image_path
      excel:   excel_path, excel_export_range, excel_driver_range
    """
    text = text.strip()

    # --- Excel: path.xlsx[driver:X,export:Y] ---
    m = EXCEL_PATH_RE.match(text)
    if m:
        excel_path  = m.group(1).strip()
        params_str  = m.group(2).strip()
        params = {}
        for part in params_str.split(","):
            part = part.strip()
            if ":" in part:
                key, val = part.split(":", 1)
                params[key.strip().lower()] = val.strip()
        export_range = params.get("export", "")
        driver_range = params.get("driver", "")
        if export_range:
            return {
                "content_type":       "excel",
                "excel_path":         excel_path,
                "excel_export_range": export_range,
                "excel_driver_range": driver_range,
            }

    # --- Picture: path ending in a supported image extension ---
    # Strip [code] and [id] tokens before checking extension
    clean = text.replace("[code]", "X").replace("[id]", "X")
    ext = os.path.splitext(clean)[1].lower()
    if ext in IMAGE_EXTENSIONS:
        return {
            "content_type": "picture",
            "image_path":   text,
        }

    # --- Chart: URL ---
    url = _extract_url_from_text(text)
    if url:
        return {
            "content_type": "chart",
            "url":          url,
            "label":        text.strip(),
        }

    return {"content_type": ""}


# ---------------------------------------------------------------------------
# Geometry helpers
# ---------------------------------------------------------------------------

def _fully_contained(inner_left, inner_top, inner_right, inner_bottom,
                     outer_left, outer_top, outer_right, outer_bottom) -> bool:
    """Return True if inner rectangle is fully within outer rectangle."""
    return (
        inner_left >= outer_left
        and inner_top >= outer_top
        and inner_right <= outer_right
        and inner_bottom <= outer_bottom
    )


# ---------------------------------------------------------------------------
# Shape classification
# ---------------------------------------------------------------------------

def _is_chart_placeholder(shape) -> bool:
    """
    Return True if this shape is a named placeholder that ChartGen
    should treat as a chart slot.

    We accept any placeholder whose name starts with "Chart" (case-insensitive),
    OR any picture placeholder, OR any object/content placeholder.
    We exclude title, body, footer, slide number, date, etc.
    """
    if not shape.is_placeholder:
        return False
    name_lower = shape.name.lower()
    # Explicit chart naming convention
    if name_lower.startswith("chart"):
        return True
    # PowerPoint placeholder types that can hold images/objects
    ph_type = shape.placeholder_format.type
    _content_types = {
        PP_PLACEHOLDER.OBJECT,
        PP_PLACEHOLDER.PICTURE,
        PP_PLACEHOLDER.CHART,
        PP_PLACEHOLDER.BITMAP,
    }
    return ph_type in _content_types


def _is_textbox(shape) -> bool:
    """Return True if shape is a free textbox (not a placeholder)."""
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    except Exception:
        return False


def _is_autoshape_with_text(shape) -> bool:
    """Return True if shape is an autoshape (rectangle etc.) with text — also counts."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE as MST
        return shape.shape_type == MST.AUTO_SHAPE and shape.has_text_frame
    except Exception:
        return False


# ---------------------------------------------------------------------------
# Core read function
# ---------------------------------------------------------------------------

def read_template(pptx_path: str) -> TemplateReadResult:
    """
    Read a .pptx template and return a TemplateReadResult.

    For each slide:
      1. Collect all chart placeholders (name, position, size)
      2. Collect all yellow textboxes that contain a URL
      3. Match yellow URL textboxes to the placeholder they sit inside
      4. Produce a cleaned copy of the pptx with yellow textboxes removed
    """
    prs = Presentation(pptx_path)
    result = TemplateReadResult(
        slide_width=int(prs.slide_width),
        slide_height=int(prs.slide_height),
    )

    # We need to track which XML elements to remove for the cleaned copy.
    # Work on the live prs first, collect removal targets, then strip.
    elements_to_remove = []   # list of (slide, sp_element)

    for slide_idx, slide in enumerate(prs.slides):
        # --- Step 1: collect chart placeholders on this slide ---
        chart_placeholders = []
        for shape in slide.shapes:
            if _is_chart_placeholder(shape):
                chart_placeholders.append(
                    PlaceholderInfo(
                        slide_index=slide_idx,
                        name=shape.name,
                        left=int(shape.left),
                        top=int(shape.top),
                        width=int(shape.width),
                        height=int(shape.height),
                    )
                )

        # --- Step 2: collect yellow textboxes on this slide ---
        yellow_shapes = []
        for shape in slide.shapes:
            if not (_is_textbox(shape) or _is_autoshape_with_text(shape)):
                continue
            rgb = _get_shape_fill_rgb(shape)
            if not _is_yellow(rgb):
                continue
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            classification = _classify_yellow_box(text)
            if not classification.get("content_type"):
                continue
            yellow_shapes.append({
                "shape":  shape,
                "left":   int(shape.left),
                "top":    int(shape.top),
                "right":  int(shape.left) + int(shape.width),
                "bottom": int(shape.top)  + int(shape.height),
                **classification,
            })

        # --- Step 3: match yellow textboxes to placeholders ---
        ph_matches: dict[str, list] = {ph.name: [] for ph in chart_placeholders}

        for yshape in yellow_shapes:
            matched = False
            for ph in chart_placeholders:
                ph_right  = ph.left + ph.width
                ph_bottom = ph.top  + ph.height
                if _fully_contained(
                    yshape["left"], yshape["top"],
                    yshape["right"], yshape["bottom"],
                    ph.left, ph.top, ph_right, ph_bottom,
                ):
                    ph_matches[ph.name].append(yshape)
                    matched = True
                    break

            if not matched:
                result.warnings.append(
                    f"Slide {slide_idx + 1}: yellow textbox not inside any "
                    f"chart placeholder — ignored. Content: "
                    f"{yshape.get('url') or yshape.get('image_path') or yshape.get('excel_path', '')}"
                )

        # --- Step 4: assign content to placeholders ---
        for ph in chart_placeholders:
            matches = ph_matches[ph.name]
            if len(matches) == 0:
                pass
            else:
                if len(matches) > 1:
                    result.warnings.append(
                        f"Slide {slide_idx + 1}, placeholder '{ph.name}': "
                        f"{len(matches)} yellow textboxes found — using first."
                    )
                m = matches[0]
                ct = m["content_type"]
                ph.content_type = ct
                if ct == "chart":
                    ph.url   = m.get("url",   "")
                    ph.label = m.get("label", "")
                elif ct == "picture":
                    ph.image_path = m.get("image_path", "")
                elif ct == "excel":
                    ph.excel_path         = m.get("excel_path",         "")
                    ph.excel_export_range = m.get("excel_export_range", "")
                    ph.excel_driver_range = m.get("excel_driver_range", "")

            result.placeholders.append(ph)

        # --- Step 5: mark all yellow textboxes on this slide for removal ---
        for yshape in yellow_shapes:
            elements_to_remove.append((slide, yshape["shape"]._element))

    # --- Step 6: produce cleaned pptx bytes ---
    import io as _io
    for slide, elem in elements_to_remove:
        try:
            elem.getparent().remove(elem)
        except Exception:
            pass  # already removed or detached

    buf = _io.BytesIO()
    prs.save(buf)
    result.cleaned_pptx_bytes = buf.getvalue()

    return result


# ---------------------------------------------------------------------------
# urls.csv helpers
# ---------------------------------------------------------------------------

def update_urls_csv(csv_path: str, new_urls: list[dict]) -> tuple[int, int]:
    """
    Add any URLs from new_urls that are not already in urls.csv.

    new_urls: list of dicts with at least {"url": str, "label": str}

    Returns (added_count, already_present_count).
    """
    import csv

    # Load existing URLs
    existing = set()
    rows = []
    if os.path.exists(csv_path):
        with open(csv_path, newline="", encoding="utf-8-sig") as f:
            reader = csv.reader(f)
            for row in reader:
                if row and row[0].strip():
                    existing.add(row[0].strip())
                    rows.append(row)

    added = 0
    already = 0
    for entry in new_urls:
        url = entry.get("url", "").strip()
        if not url:
            continue
        if url in existing:
            already += 1
        else:
            label = entry.get("label", "").strip()
            rows.append([url, label])
            existing.add(url)
            added += 1

    if added > 0:
        with open(csv_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            for row in rows:
                writer.writerow(row)

    return added, already
