"""
assembly_engine.py
Executes a Running Order against a PowerPoint template to produce one output file.
Only this module touches python-pptx directly.
"""

import os
import io
import shutil
import time
import traceback

from pptx import Presentation
from pptx.util import Emu
from PIL import Image as PILImage

from modules.m05_chart_engine.cache_reader import load_shape
from modules.m05_chart_engine.base_charts import render_chart
from modules.m12_local_config.local_config import ReportContext, build_report_context
from modules.m04_data_shapes.shapes import PopulationShape, filter_shape
from modules.m07_insert_picture.insert_picture import insert_picture
from modules.m08_insert_from_excel.insert_from_excel import (
    open_excel, close_excel, insert_from_excel
)


def build_population_shapes(data_shape, populations_str: str,
                             units: list, report_context) -> list:
    """
    Build an ordered list of PopulationShape objects from a '^'-delimited
    populations string. The first token defines the scope (the full set being
    compared); each subsequent token is an independent subset of that scope.
    Tokens: 'All' (every unit), 'Selected' (current organisation),
    'Name()' (selected unit's own group), 'Name(Value)' (the named group).
    Returns [] if populations_str is blank or the scope resolves empty.
    Unit ids are compared as strings throughout.
    """
    if not populations_str or not populations_str.strip():
        return []

    shape_ids = {u.unit_id for u in _get_shape_units(data_shape)}
    if not shape_ids:
        return []

    unit_lookup = {r["unit_id"]: r for r in units}
    selected_id = report_context.unit_id if report_context else None
    selected_org_id = str(report_context.organisation_id) if report_context else None

    def _resolve(token: str, scope_ids: set):
        """Resolve one token to (unit_ids, label) within scope_ids, or None."""
        if token == "All":
            return set(scope_ids), "All"

        if token == "Selected":
            if not selected_org_id:
                return None
            ids = {
                r["unit_id"] for r in units
                if str(r["organisation_id"]) == selected_org_id
                and r["unit_id"] in scope_ids
            }
            return ids, "Selected"

        if token.endswith(")") and "(" in token:
            col_name, value = token[:-1].split("(", 1)
            col = col_name + "()"
            if not value:  # Name() — selected unit's own group
                if not selected_id or selected_id not in unit_lookup:
                    return None
                value = unit_lookup[selected_id].get(col, "")
                if not value or value == "x":  # 'x' marks deliberate absence, same as blank
                    return None
            ids = {
                r["unit_id"] for r in units
                if r.get(col) == value
                and r["unit_id"] in scope_ids
            }
            return ids, value

        return None

    results = []
    scope_ids = set(shape_ids)

    for i, token in enumerate(t.strip() for t in populations_str.split("^")):
        if not token:
            continue

        resolved = _resolve(token, scope_ids)
        if resolved is None:
            if i == 0:
                return []
            continue
        token_ids, label = resolved

        if i == 0:
            if not token_ids:
                return []
            scope_ids = token_ids

        filtered = filter_shape(data_shape, token_ids)
        results.append(PopulationShape(role=token, label=label, shape=filtered))

    return results


def _get_shape_units(data_shape) -> list:
    """Return the flat list of units from any shape type."""
    from modules.m04_data_shapes.shapes import (
        NumericSeries, NumericCompositional, CategoricalCompositional
    )
    if isinstance(data_shape, NumericSeries):
        return data_shape.units
    elif isinstance(data_shape, (NumericCompositional, CategoricalCompositional)):
        return data_shape.metrics[0].units if data_shape.metrics else []
    return []


# ---------------------------------------------------------------------------
# Assembly context — passed through every function call in a run
# ---------------------------------------------------------------------------

class AssemblyContext:
    def __init__(self):
        self.prs: Presentation = None
        self.output_path: str = ""
        self.template_path: str = ""
        self.log: list[dict] = []
        self.autotable_stats: dict = {}
        self.report_context = None      # set by run_running_order
        self.default_populations: str = ""  # set by set_default_populations row


# ---------------------------------------------------------------------------
# Running Order function implementations
# ---------------------------------------------------------------------------

def create_ppt(ctx: AssemblyContext, row: dict, settings: dict) -> dict:
    """
    Open the cleaned template and set the output path.
    Does not write to disk — save_ppt/save_pdf do that.
    """
    template_path = settings.get("cleaned_template_path", "").strip()
    if not template_path or not os.path.exists(template_path):
        # Fall back to the original template if no cleaned version exists
        template_path = settings.get("ppt_template_path", "").strip()

    if not template_path or not os.path.exists(template_path):
        return _err(row, "create_ppt: no template found. Check settings.")

    output_folder = _ensure_output_folder(settings)
    unit_name = (settings.get("reporting_unit_name") or "").strip()
    if not unit_name:
        unit_name = str(settings.get("selected_unit_id") or "output").strip()
    safe_name = _safe_filename(unit_name)
    output_path = os.path.join(output_folder, "pptx", f"{safe_name}.pptx")

    ctx.prs = Presentation(template_path)
    ctx.output_path = output_path
    ctx.template_path = template_path

    return _ok(row, f"Template opened: {os.path.basename(template_path)}")


def set_default_populations(ctx: AssemblyContext, row: dict, settings: dict) -> dict:
    """
    Store the default populations string on AssemblyContext.
    Subsequent insert_chart rows inherit it unless overridden per row.
    """
    populations = str(row.get("populations", "") or "").strip()
    ctx.default_populations = populations
    return _ok(row, f"Default populations set: '{populations}'")


def insert_chart(ctx: AssemblyContext, row: dict, settings: dict) -> dict:
    """Render a Base Chart from cached data and insert it at the placeholder position."""
    if ctx.prs is None:
        return _err(row, "insert_chart: no open presentation (create_ppt not called?).")

    cache_file = str(row.get("cache_file") or "").strip()
    if cache_file.lower() == "none":
        cache_file = ""
    chart_type_ref = str(row.get("chart_type_ref", "")).strip()
    slide_index = _int_or_none(row.get("slide_index"))
    placeholder = str(row.get("placeholder", "")).strip()

    # Position / size from the Running Order row (written from template at generation time)
    left_emu = _int_or_none(row.get("left_emu"))
    top_emu = _int_or_none(row.get("top_emu"))
    width_emu = _int_or_none(row.get("width_emu"))
    height_emu = _int_or_none(row.get("height_emu"))

    # Validate required fields
    missing = []
    if not cache_file:  missing.append("cache_file")
    if not chart_type_ref: missing.append("chart_type_ref")
    if slide_index is None: missing.append("slide_index")
    if None in (left_emu, top_emu, width_emu, height_emu):
        missing.append("position/size EMU values")
    if missing:
        return _err(row, f"insert_chart: missing required fields: {', '.join(missing)}")

    # --- Resolve populations for this chart ---
    row_populations = str(row.get("populations", "") or "").strip()
    populations_str = row_populations if row_populations else ctx.default_populations

    render_context = ctx.report_context

    # --- Load data shape ---
    try:
        data_shape, shape_type = _load_chart_data(cache_file, settings.get("workfile_state"))
    except Exception as e:
        return _err(row, f"insert_chart: failed to load cache '{cache_file}': {e}")

    # --- Build population shapes ---
    population_shapes = []
    if render_context is not None and populations_str:
        workfile_state = settings.get("workfile_state")
        units = workfile_state.units
        try:
            population_shapes = build_population_shapes(
                data_shape, populations_str, units, render_context
            )
        except Exception as e:
            return _err(row, f"insert_chart: failed to build population shapes: {e}")

    # Fall back to full unfiltered shape if no populations resolved
    if not population_shapes:
        from modules.m04_data_shapes.shapes import PopulationShape
        population_shapes = [PopulationShape(role="All", label="All", shape=data_shape)]

    # --- Render chart image ---
    try:
        image_bytes, autotable_stats = _render_chart_image(
            chart_type_ref, population_shapes, width_emu, height_emu, render_context
        )
    except Exception as e:
        return _err(row, f"insert_chart: render failed for '{chart_type_ref}': {e}")

    # Store autotable stats keyed by placeholder name
    ctx.autotable_stats[placeholder] = autotable_stats

    # --- Insert into slide ---
    try:
        _insert_image_at_position(
            ctx.prs, slide_index,
            image_bytes, left_emu, top_emu, width_emu, height_emu
        )
    except Exception as e:
        return _err(row, f"insert_chart: failed to insert image on slide {slide_index}: {e}")

    return _ok(row, f"Chart '{chart_type_ref}' inserted at '{placeholder}' (slide {slide_index + 1})")


def update_text(ctx: AssemblyContext, row: dict, settings: dict) -> dict:
    """
    Replace text tags in the presentation with values for the current reporting unit.
    Tags: [selected-reporting-unit-name] → unit_name.
    """
    if ctx.prs is None:
        return _err(row, "update_text: no open presentation (create_ppt not called?).")

    rc = ctx.report_context
    tokens = {}
    if rc:
        tokens["[selected-reporting-unit-name]"] = rc.unit_name or ""

    if not tokens:
        return _ok(row, "update_text: no tags to replace (no ReportContext).")

    # XML namespace for DrawingML text runs
    _nsmap = "a"
    _run_tag = "{http://schemas.openxmlformats.org/drawingml/2006/main}r"
    _t_tag   = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"

    replacements = 0

    for slide in ctx.prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                runs = para.runs
                if not runs:
                    continue

                # Check whether the full paragraph text contains any tag
                full_text = "".join(r.text for r in runs)
                needs_replace = any(tok in full_text for tok in tokens)
                if not needs_replace:
                    continue

                # Apply all tag replacements to the full paragraph text
                replaced = full_text
                for token, value in tokens.items():
                    if token in replaced:
                        replaced = replaced.replace(token, value)
                        replacements += 1

                # Write the replaced text into the first run's <a:t> element,
                # then delete all subsequent runs from the paragraph XML.
                first_run_xml = runs[0]._r        # lxml element for the run
                t_elem = first_run_xml.find(_t_tag)
                if t_elem is not None:
                    t_elem.text = replaced

                # Remove every run element after the first from the paragraph XML
                para_xml = para._p
                run_elements = para_xml.findall(_run_tag)
                for run_elem in run_elements[1:]:
                    para_xml.remove(run_elem)

    return _ok(row, f"update_text: {replacements} replacement(s) made.")


def empty_placeholder(ctx: AssemblyContext, row: dict, settings: dict) -> dict:
    """No-op. Placeholder has no content assigned."""
    ph = row.get("placeholder", "")
    return _ok(row, f"empty_placeholder: '{ph}' skipped (no content assigned)")


def save_ppt(ctx: AssemblyContext, row: dict, settings: dict) -> dict:
    """Save the completed output as a .pptx file."""
    if ctx.prs is None:
        return _err(row, "save_ppt: no open presentation.")
    try:
        os.makedirs(os.path.dirname(ctx.output_path), exist_ok=True)
        ctx.prs.save(ctx.output_path)
        return _ok(row, f"Saved: {ctx.output_path}")
    except Exception as e:
        return _err(row, f"save_ppt: {e}")


def save_pdf(ctx: AssemblyContext, row: dict, settings: dict) -> dict:
    """
    Save the completed output as a .pdf using COM automation (Windows/PowerPoint only).
    Falls back gracefully on non-Windows or if PowerPoint is not available.
    """
    if ctx.prs is None:
        return _err(row, "save_pdf: no open presentation.")

    pdf_dir = os.path.join(os.path.dirname(os.path.dirname(ctx.output_path)), "pdf")
    pdf_path = os.path.join(pdf_dir, os.path.basename(ctx.output_path).replace(".pptx", ".pdf"))
    os.makedirs(pdf_dir, exist_ok=True)

    # Ensure the pptx is saved first (COM needs a file on disk to open)
    try:
        ctx.prs.save(ctx.output_path)
    except Exception as e:
        return _err(row, f"save_pdf: could not save .pptx before PDF export: {e}")

    try:
        import comtypes.client
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        deck = powerpoint.Presentations.Open(os.path.abspath(ctx.output_path))
        deck.ExportAsFixedFormat(
                os.path.abspath(pdf_path),
                2,   # ppFixedFormatTypePDF
                2,   # Intent: ppFixedFormatIntentPrint (vs 1 = screen quality)
            )
        deck.Close()
        powerpoint.Quit()
        return _ok(row, f"PDF saved: {pdf_path}")
    except ImportError:
        return _err(row, "save_pdf: comtypes not available — PDF export requires Windows + PowerPoint.")
    except Exception as e:
        return _err(row, f"save_pdf: COM export failed: {e}")


# ---------------------------------------------------------------------------
# Dispatch map  —  function name -> callable
# ---------------------------------------------------------------------------

FUNCTION_MAP = {
    "create_ppt":               create_ppt,
    "set_default_populations":  set_default_populations,
    "insert_chart":             insert_chart,
    "insert_picture":           insert_picture,
    "insert_from_excel":        insert_from_excel,
    "open_excel":               open_excel,
    "close_excel":              close_excel,
    "update_text":              update_text,
    "empty_placeholder":        empty_placeholder,
    "save_ppt":                 save_ppt,
    "save_pdf":                 save_pdf,
}


# ---------------------------------------------------------------------------
# Run a complete Running Order
# ---------------------------------------------------------------------------

def run_running_order(rows: list[dict], settings: dict,
                      on_progress=None, ctx: AssemblyContext = None) -> dict:
    """
    Execute a list of Running Order rows (already filtered to enabled only).

    settings dict must contain at minimum:
      ppt_template_path, cleaned_template_path, workfile_folder,
      reporting_unit_name, workfile_state

    Returns:
    {"status": "ok" | "error", "output_path": str, "elapsed": float, "log": list[dict]}
    """
    # Use a shared context if provided (batch run), otherwise create a fresh one.
    if ctx is None:
        ctx = AssemblyContext()

    workfile_state = settings.get("workfile_state")
    units = workfile_state.units
    ctx.report_context = build_report_context(settings, units)

    t_start = time.perf_counter()

    normal_rows = [r for r in rows if str(r.get("scope", "normal")).strip() == "normal"]
    rows_to_run = normal_rows

    total = len(rows_to_run)

    for i, row in enumerate(rows_to_run):
        func_name = str(row.get("function", "")).strip()

        if on_progress:
            on_progress(i + 1, total, f"Row {row.get('row_id', i+1)}: {func_name}")

        func = FUNCTION_MAP.get(func_name)
        if func is None:
            result = _err(row, f"Unknown function: '{func_name}'")
        else:
            try:
                result = func(ctx, row, settings)
            except Exception as e:
                result = _err(row, f"Unhandled exception in '{func_name}': {traceback.format_exc()}")

        ctx.log.append(result)

        # Abort on error in structural functions
        if result["status"] == "error" and func_name in ("create_ppt",):
            ctx.log.append({"status": "aborted",
                            "message": "Batch aborted after create_ppt failure."})
            break

    elapsed = time.perf_counter() - t_start
    overall_status = "ok" if all(r["status"] in ("ok", "skip") for r in ctx.log) else "error"

    return {
        "status": overall_status,
        "output_path": ctx.output_path,
        "elapsed": elapsed,
        "log": ctx.log,
    }


# ---------------------------------------------------------------------------
# Private helpers — internal to insert_chart sub-steps
# ---------------------------------------------------------------------------

def _load_chart_data(cache_file: str, workfile_state=None):
    """Load a canonical data shape from the cache. Sub-step of insert_chart."""
    return load_shape(cache_file, workfile_state)


def _render_chart_image(chart_type_ref: str, population_shapes: list, width_emu: int, height_emu: int,
                        report_context=None):
    """
    Render a Matplotlib chart to PNG bytes sized to the placeholder.
    Sub-step of insert_chart.
    """
    NARROWER_EMU = 6858000
    width_pct  = max(10, int(min(100, (width_emu  / NARROWER_EMU) * 100)))
    height_pct = max(10, int(min(100, (height_emu / NARROWER_EMU) * 100)))

    image_bytes, autotable_stats = render_chart(
        chart_type_ref, population_shapes,
        width=width_pct, height=height_pct,
        report_context=report_context,
    )
    return image_bytes, autotable_stats


def _insert_image_at_position(prs: Presentation, slide_index: int,
                               image_bytes, left_emu: int, top_emu: int,
                               width_emu: int, height_emu: int):
    """
    Insert a PNG image at the exact EMU position on the given slide.
    Sub-step of insert_chart.
    """
    if slide_index >= len(prs.slides):
        raise IndexError(
            f"Slide index {slide_index} out of range "
            f"(template has {len(prs.slides)} slides)."
        )
    slide = prs.slides[slide_index]
    image_bytes.seek(0)
    slide.shapes.add_picture(
        image_bytes,
        Emu(left_emu), Emu(top_emu),
        Emu(width_emu), Emu(height_emu),
    )


# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------

def _ensure_output_folder(settings: dict) -> str:
    output_folder = settings.get("outputs_folder", "").strip()
    if not output_folder:
        workfile_folder = settings.get("workfile_folder", "").strip()
        output_folder = os.path.join(workfile_folder, "outputs") if workfile_folder else "outputs"
    os.makedirs(output_folder, exist_ok=True)
    return output_folder


def _safe_filename(name: str) -> str:
    """Strip characters that are unsafe in filenames."""
    import re
    safe = re.sub(r'[\\/:*?"<>|]', "_", name)
    return safe.strip("_") or "output"


def _int_or_none(value):
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _ok(row: dict, message: str) -> dict:
    return {
        "status": "ok",
        "row_id": row.get("row_id", ""),
        "function": row.get("function", ""),
        "message": message,
    }


def _err(row: dict, message: str) -> dict:
    return {
        "status": "error",
        "row_id": row.get("row_id", ""),
        "function": row.get("function", ""),
        "message": message,
    }
