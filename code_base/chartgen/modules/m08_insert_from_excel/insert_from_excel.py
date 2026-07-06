"""
insert_from_excel.py
Running Order functions: open_excel, close_excel, insert_from_excel. Captures a named Excel
range as an image and inserts it into the PowerPoint output.

Row fields:
  open_excel / close_excel:   excel_path
  insert_from_excel:          excel_path, export_range, driver_range (optional),
                               left_emu, top_emu, width_emu, slide_index
"""

import os
import io
import tempfile
from pptx.util import Emu


# ---------------------------------------------------------------------------
# AssemblyContext extension
# ---------------------------------------------------------------------------
# open_excel stores its COM object on ctx.excel_workbooks — a dict keyed by
# normalised file path. This supports multiple open workbooks in one batch.
# close_excel removes the entry and releases the COM object.


def _normalise_path(path: str) -> str:
    """Return a normalised absolute path for use as a dict key."""
    return os.path.normpath(os.path.abspath(path.strip()))


# ---------------------------------------------------------------------------
# open_excel
# ---------------------------------------------------------------------------

def open_excel(ctx, row: dict, settings: dict) -> dict:
    """
    Open an Excel workbook via COM and hold the reference on AssemblyContext.
    Scope: batch_open — runs once at the start of a batch.
    """
    from modules.m06_assembly_engine.assembly_engine import _err, _ok

    excel_path = str(row.get("excel_path", "")).strip()
    if not excel_path:
        return _err(row, "open_excel: no excel_path specified.")
    if not os.path.exists(excel_path):
        return _err(row, f"open_excel: file not found: '{excel_path}'")

    norm_path = _normalise_path(excel_path)

    # Initialise workbook store on context if not present
    if not hasattr(ctx, "excel_workbooks"):
        ctx.excel_workbooks = {}

    if norm_path in ctx.excel_workbooks:
        return _ok(row, f"open_excel: '{os.path.basename(excel_path)}' already open.")

    try:
        import pythoncom
        import win32com.client as win32
        pythoncom.CoInitialize()
    except ImportError:
        return _err(row, "open_excel: pywin32 not available. "
                         "Install with: pip install pywin32")

    try:
        xl = win32.Dispatch("Excel.Application")
        xl.Visible = True          # visible during testing — set False once confirmed working
        xl.DisplayAlerts = False

        wb = xl.Workbooks.Open(os.path.abspath(excel_path))

        # Disable auto-calculation — each insert_from_excel will Calculate() explicitly
        xl.Calculation = -4135  # xlCalculationManual

        # Confirm workbook is genuinely open and accessible
        wb_name = wb.Name
        sheet_count = wb.Sheets.Count
        if sheet_count == 0:
            raise RuntimeError("Workbook opened but contains no sheets.")

        ctx.excel_workbooks[norm_path] = {
            "app": xl,
            "workbook": wb,
        }

        return _ok(row, f"open_excel: opened '{wb_name}' ({sheet_count} sheets).")

    except Exception as e:
        return _err(row, f"open_excel: failed to open workbook: {e}")


# ---------------------------------------------------------------------------
# close_excel
# ---------------------------------------------------------------------------

def close_excel(ctx, row: dict, settings: dict) -> dict:
    """
    Close a workbook and release the COM object.

    Scope: batch_close — runs once after the last report in a batch.
    """
    from modules.m06_assembly_engine.assembly_engine import _err, _ok

    excel_path = str(row.get("excel_path", "")).strip()
    if not excel_path:
        return _err(row, "close_excel: no excel_path specified.")

    norm_path = _normalise_path(excel_path)

    if not hasattr(ctx, "excel_workbooks") or norm_path not in ctx.excel_workbooks:
        return _ok(row, f"close_excel: '{os.path.basename(excel_path)}' "
                        f"was not open — nothing to close.")

    entry = ctx.excel_workbooks.pop(norm_path)
    try:
        entry["workbook"].Close(SaveChanges=False)
        entry["app"].Quit()
        import pythoncom
        pythoncom.CoUninitialize()
    except Exception as e:
        return _err(row, f"close_excel: error closing workbook: {e}")

    return _ok(row, f"close_excel: closed '{os.path.basename(excel_path)}'.")


# ---------------------------------------------------------------------------
# insert_from_excel
# ---------------------------------------------------------------------------

def insert_from_excel(ctx, row: dict, settings: dict) -> dict:
    """
    Write unit_id to the driver range (if specified), force workbook
    recalculation, capture the export range as a PNG, and insert into the
    presentation at the placeholder position.

    Scope: normal — runs once per report.
    """
    from modules.m06_assembly_engine.assembly_engine import _err, _ok

    if ctx.prs is None:
        return _err(row, "insert_from_excel: no open presentation "
                         "(create_ppt must run first).")

    if ctx.report_context is None:
        return _err(row, "insert_from_excel: no ReportContext available.")

    # --- Resolve row fields ---
    excel_path   = str(row.get("excel_path",   "")).strip()
    export_range = str(row.get("export_range", "")).strip()
    driver_range = str(row.get("driver_range", "")).strip()

    if not excel_path:
        return _err(row, "insert_from_excel: no excel_path specified.")
    if not export_range:
        return _err(row, "insert_from_excel: no export_range specified.")

    norm_path = _normalise_path(excel_path)

    if not hasattr(ctx, "excel_workbooks") or norm_path not in ctx.excel_workbooks:
        return _err(row, f"insert_from_excel: workbook '{os.path.basename(excel_path)}' "
                         f"is not open. Add an open_excel row (scope=batch_open) "
                         f"before this row in the Running Order.")

    entry = ctx.excel_workbooks[norm_path]
    wb    = entry["workbook"]
    xl    = entry["app"]

    # --- Write driver cell ---
    if driver_range:
        try:
            driver_cell = wb.Names(driver_range).RefersToRange
            driver_cell.Value = ctx.report_context.unit_id
        except Exception as e:
            return _err(row, f"insert_from_excel: could not write to driver range "
                             f"'{driver_range}': {e}")

    # --- Force recalculation ---
    try:
        xl.Calculate()
    except Exception as e:
        import traceback
        detail = f"insert_from_excel: workbook calculation failed\nType: {type(e).__name__}\nRepr: {e!r}\n\n{traceback.format_exc()}"
        # Write to log file alongside the workfile folder for inspection
        try:
            log_path = os.path.join(settings.get("workfile_folder", "."), "chartgen_error.log")
            with open(log_path, "w", encoding="utf-8") as f:
                f.write(detail)
        except Exception:
            pass
        return _err(row, f"insert_from_excel: workbook calculation failed — see chartgen_error.log in workfile folder")

    # --- Capture export range as image ---
    try:
        export_rng = wb.Names(export_range).RefersToRange
    except Exception as e:
        return _err(row, f"insert_from_excel: could not find export range "
                         f"'{export_range}': {e}")

    try:
        image_bytes = _capture_range_as_png(export_rng, xl)
    except Exception as e:
        return _err(row, f"insert_from_excel: image capture failed: {e}")

    # --- Resolve position ---
    try:
        left_emu  = int(row.get("left_emu",  0))
        top_emu   = int(row.get("top_emu",   0))
        width_emu = int(row.get("width_emu", 0))
        slide_idx = int(row.get("slide_index", 0))
    except (ValueError, TypeError) as e:
        return _err(row, f"insert_from_excel: invalid position values: {e}")

    if width_emu <= 0:
        return _err(row, "insert_from_excel: width_emu must be greater than zero.")

    # --- Calculate height from aspect ratio ---
    try:
        from PIL import Image as PILImage
        with PILImage.open(io.BytesIO(image_bytes)) as img:
            img_w, img_h = img.size
        if img_w <= 0:
            return _err(row, "insert_from_excel: captured image has zero width.")
        height_emu = int(width_emu * (img_h / img_w))
    except Exception as e:
        return _err(row, f"insert_from_excel: could not determine image dimensions: {e}")

    # --- Insert into slide ---
    try:
        if slide_idx >= len(ctx.prs.slides):
            return _err(row, f"insert_from_excel: slide index {slide_idx} out of range "
                             f"(presentation has {len(ctx.prs.slides)} slides).")
        slide = ctx.prs.slides[slide_idx]
        slide.shapes.add_picture(
            io.BytesIO(image_bytes),
            Emu(left_emu), Emu(top_emu),
            Emu(width_emu), Emu(height_emu),
        )
    except Exception as e:
        return _err(row, f"insert_from_excel: failed to insert image into slide: {e}")

    return _ok(row, f"insert_from_excel: inserted '{export_range}' from "
                    f"'{os.path.basename(excel_path)}' at slide {slide_idx + 1}.")


# ---------------------------------------------------------------------------
# Image capture helper
# ---------------------------------------------------------------------------

def _capture_range_as_png(excel_range, xl_app) -> bytes:
    """Copy an Excel range to the clipboard as a picture and return PNG bytes."""
    import time
    from PIL import ImageGrab

    # CopyPicture: Appearance=1 (xlScreen), Format=2 (xlBitmap)
    excel_range.CopyPicture(Appearance=1, Format=2)

    # Brief pause to allow clipboard to be populated
    time.sleep(0.3)

    img = ImageGrab.grabclipboard()
    if img is None:
        raise RuntimeError(
            "Clipboard image was None after CopyPicture. "
            "Ensure Excel is not minimised and the named range is visible."
        )

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()
