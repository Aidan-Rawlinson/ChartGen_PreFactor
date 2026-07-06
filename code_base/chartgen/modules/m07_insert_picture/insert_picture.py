"""
insert_picture.py
Running Order function: inserts a static image at a named placeholder, with [code]/[id] token
substitution. Height is derived from the image's aspect ratio, not the placeholder.
"""

import os
import io
from PIL import Image as PILImage
from pptx.util import Emu


SUPPORTED_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".wmf", ".emf"}


def _substitute_tokens(path: str, report_context) -> str:
    """Replace [code] and [id] tokens in a path string with values from ReportContext."""
    if report_context is None:
        return path
    path = path.replace("[code]", str(report_context.unit_code))
    path = path.replace("[id]",   str(report_context.unit_id))
    return path


def insert_picture(ctx, row: dict, settings: dict) -> dict:
    """
    Insert a static image at the placeholder position defined in the Running Order row.

    Required row fields: image_path, left_emu, top_emu, width_emu, slide_index.
    """
    from modules.m06_assembly_engine.assembly_engine import _err, _ok

    if ctx.prs is None:
        return _err(row, "insert_picture: no open presentation (create_ppt must run first).")

    # --- Resolve path ---
    raw_path = str(row.get("image_path", "")).strip()
    if not raw_path:
        return _err(row, "insert_picture: no image_path specified in Running Order row.")

    resolved_path = _substitute_tokens(raw_path, ctx.report_context)

    ext = os.path.splitext(resolved_path)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return _err(row, f"insert_picture: unsupported file type '{ext}'. "
                         f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")

    if not os.path.exists(resolved_path):
        return _err(row, f"insert_picture: file not found: '{resolved_path}'")

    # --- Read dimensions ---
    try:
        left_emu  = int(row.get("left_emu",  0))
        top_emu   = int(row.get("top_emu",   0))
        width_emu = int(row.get("width_emu", 0))
        slide_idx = int(row.get("slide_index", 0))
    except (ValueError, TypeError) as e:
        return _err(row, f"insert_picture: invalid position values in row: {e}")

    if width_emu <= 0:
        return _err(row, "insert_picture: width_emu must be greater than zero.")

    # --- Calculate height from aspect ratio ---
    try:
        with PILImage.open(resolved_path) as img:
            img_w, img_h = img.size
        if img_w <= 0:
            return _err(row, "insert_picture: image has zero width.")
        height_emu = int(width_emu * (img_h / img_w))
    except Exception as e:
        return _err(row, f"insert_picture: could not read image dimensions: {e}")

    # --- Insert into slide ---
    try:
        if slide_idx >= len(ctx.prs.slides):
            return _err(row, f"insert_picture: slide index {slide_idx} out of range "
                             f"(presentation has {len(ctx.prs.slides)} slides).")
        slide = ctx.prs.slides[slide_idx]
        slide.shapes.add_picture(
            resolved_path,
            Emu(left_emu), Emu(top_emu),
            Emu(width_emu), Emu(height_emu),
        )
    except Exception as e:
        return _err(row, f"insert_picture: failed to insert image: {e}")

    return _ok(row, f"insert_picture: inserted '{os.path.basename(resolved_path)}' "
                    f"at slide {slide_idx + 1} ({width_emu} × {height_emu} EMU).")
